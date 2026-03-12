"""Service for report template parsing, fill, and export."""

from __future__ import annotations

import io
import json
import logging
import re
from dataclasses import asdict
from pathlib import Path
from uuid import uuid4

from openpyxl import load_workbook

from app.domain.entities.report import FillableRegion, ReportJob, ReportTemplate

MARKER_PATTERN = re.compile(r"\{\{(\w+)\}\}")

logger = logging.getLogger(__name__)


class ReportService:
    """Manages report templates and fill jobs."""

    def __init__(
        self,
        template_repo,
        job_repo,
        file_storage,
        connector_service=None,
        openai_api_key: str = "",
        openai_model: str = "gpt-4o",
    ) -> None:
        self._template_repo = template_repo
        self._job_repo = job_repo
        self._file_storage = file_storage
        self._connector_service = connector_service
        self._openai_api_key = openai_api_key
        self._openai_model = openai_model

        # Lazily create OpenAI client only when needed
        self._openai_client = None
        if openai_api_key:
            from openai import AsyncOpenAI
            self._openai_client = AsyncOpenAI(api_key=openai_api_key)

    # ------------------------------------------------------------------
    # Templates
    # ------------------------------------------------------------------

    async def upload_template(self, filename: str, file_bytes: bytes) -> ReportTemplate:
        ext = Path(filename).suffix.lower().lstrip(".")
        if ext not in ("xlsx", "pptx"):
            raise ValueError(f"Unsupported template format: {ext}")

        storage_path = f"report_templates/{uuid4()}/{filename}"
        await self._file_storage.store(file_bytes, storage_path)

        if ext == "xlsx":
            regions = _detect_xlsx_regions(file_bytes)
        else:
            regions = _detect_pptx_regions(file_bytes)

        template = ReportTemplate(
            name=filename,
            file_format=ext,
            file_path=storage_path,
            regions=[asdict(r) for r in regions],
        )
        return await self._template_repo.create(template)

    async def get_template(self, template_id: str) -> ReportTemplate | None:
        return await self._template_repo.get_by_id(template_id)

    async def list_templates(self) -> list[ReportTemplate]:
        return await self._template_repo.list_all()

    # ------------------------------------------------------------------
    # Jobs
    # ------------------------------------------------------------------

    async def create_job(self, template_id: str, name: str) -> ReportJob:
        template = await self._template_repo.get_by_id(template_id)
        if template is None:
            raise ValueError(f"Template {template_id} not found")
        job = ReportJob(template_id=template_id, name=name)
        return await self._job_repo.create(job)

    async def update_fills(self, job_id: str, fills: dict) -> ReportJob:
        job = await self._job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Job {job_id} not found")
        job.fills = {**job.fills, **fills}
        return await self._job_repo.update(job)

    async def generate(self, job_id: str) -> ReportJob:
        job = await self._job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Job {job_id} not found")
        template = await self._template_repo.get_by_id(job.template_id)
        if template is None:
            raise ValueError(f"Template {job.template_id} not found")

        # Read template file
        template_path = await self._file_storage.retrieve(template.file_path)
        template_bytes = Path(template_path).read_bytes()

        # Apply fills
        if template.file_format == "xlsx":
            output_bytes = _fill_xlsx(template_bytes, job.fills, template.regions)
        else:
            output_bytes = _fill_pptx(template_bytes, job.fills, template.regions)

        # Store output
        ext = template.file_format
        output_storage_path = f"report_outputs/{job.id}/output.{ext}"
        await self._file_storage.store(output_bytes, output_storage_path)

        job.output_file_path = output_storage_path
        job.status = "completed"
        return await self._job_repo.update(job)

    async def get_job(self, job_id: str) -> ReportJob | None:
        return await self._job_repo.get_by_id(job_id)

    async def list_jobs(self) -> list[ReportJob]:
        return await self._job_repo.list_all()

    async def download(self, job_id: str) -> tuple[bytes, str]:
        job = await self._job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Job {job_id} not found")
        if not job.output_file_path:
            raise ValueError("Job output has not been generated yet")
        template = await self._template_repo.get_by_id(job.template_id)
        file_path = await self._file_storage.retrieve(job.output_file_path)
        file_bytes = Path(file_path).read_bytes()
        filename = f"{job.name}.{template.file_format}" if template else Path(job.output_file_path).name
        return file_bytes, filename

    # ------------------------------------------------------------------
    # AI Fill
    # ------------------------------------------------------------------

    async def ai_fill(
        self,
        job_id: str,
        connectors: list[str],
        prompt: str | None = None,
    ) -> ReportJob:
        """Use AI to auto-fill template regions from connected file context."""
        job = await self._job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Job {job_id} not found")

        template = await self._template_repo.get_by_id(job.template_id)
        if template is None:
            raise ValueError(f"Template {job.template_id} not found")

        if not self._openai_client:
            raise ValueError("OpenAI API key is not configured")

        regions = template.regions
        if not regions:
            return job

        # Gather file context from connected sources for each region
        file_context_parts: list[str] = []
        if self._connector_service and connectors:
            for region in regions:
                label = region.get("label", "")
                headers = region.get("headers", [])
                query = f"{label} {' '.join(headers)}".strip()
                if not query:
                    continue
                try:
                    results = await self._connector_service.search_files(query)
                    for f in results:
                        text = f.text_content or ""
                        if text:
                            file_context_parts.append(
                                f"--- File: {f.name} ---\n{text}"
                            )
                except Exception:
                    logger.warning(
                        "Failed to search files for region %s", label, exc_info=True
                    )

        # Build the region structure description
        region_descriptions: list[str] = []
        for region in regions:
            rid = region.get("region_id", "")
            label = region.get("label", "")
            headers = region.get("headers", [])
            row_count = region.get("row_count", 0)
            region_descriptions.append(
                f'- region_id: "{rid}", label: "{label}", '
                f'headers: {json.dumps(headers)}, row_count: {row_count}'
            )

        system_message = (
            "You are a report-filling assistant. Given a template structure and "
            "optional file context, generate data to fill the template regions.\n\n"
            "Return a JSON object where each key is a region_id and each value has "
            'a "rows" array. Each row is an array of string values matching the '
            "headers in order. Generate the number of rows indicated by row_count "
            "(or fewer if data is insufficient).\n\n"
            "Example output:\n"
            '{\n  "region-abc": {\n    "rows": [["val1", "val2"], ["val3", "val4"]]\n  }\n}'
        )

        user_parts: list[str] = []
        user_parts.append("## Template Regions\n" + "\n".join(region_descriptions))

        if file_context_parts:
            # Deduplicate and limit context size
            seen: set[str] = set()
            unique_parts: list[str] = []
            for part in file_context_parts:
                if part not in seen:
                    seen.add(part)
                    unique_parts.append(part)
            file_context = "\n\n".join(unique_parts)
            # Truncate to ~30k chars to stay within context limits
            if len(file_context) > 30000:
                file_context = file_context[:30000] + "\n... (truncated)"
            user_parts.append("## File Context\n" + file_context)

        if prompt:
            user_parts.append("## User Instructions\n" + prompt)

        user_message = "\n\n".join(user_parts)

        response = await self._openai_client.chat.completions.create(
            model=self._openai_model,
            messages=[
                {"role": "system", "content": system_message},
                {"role": "user", "content": user_message},
            ],
            response_format={"type": "json_object"},
            temperature=0.2,
        )

        raw = response.choices[0].message.content or "{}"
        try:
            ai_fills = json.loads(raw)
        except json.JSONDecodeError:
            logger.error("Failed to parse AI fill response: %s", raw)
            ai_fills = {}

        # Convert AI fills into the job fills format
        # Store by both region_id (for frontend lookup) and label (for _fill_xlsx)
        for region in regions:
            rid = region.get("region_id", "")
            label = region.get("label", "")
            if rid in ai_fills and "rows" in ai_fills[rid]:
                rows = ai_fills[rid]["rows"]
                job.fills[rid] = {"rows": rows}  # for frontend to read by region_id
                job.fills[label] = rows  # for _fill_xlsx to read by sheet name

        return await self._job_repo.update(job)


# ======================================================================
# Region detection helpers
# ======================================================================


def _detect_xlsx_regions(file_bytes: bytes) -> list[FillableRegion]:
    """Scan an XLSX workbook for tables containing {{marker}} placeholders."""
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    regions: list[FillableRegion] = []
    for ws in wb.worksheets:
        # Row 1 = headers
        headers: list[str] = []
        for cell in ws[1]:
            headers.append(str(cell.value) if cell.value is not None else "")
        # Check remaining rows for markers
        markers_found: set[str] = set()
        row_count = 0
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                val = str(cell.value) if cell.value is not None else ""
                markers_found.update(MARKER_PATTERN.findall(val))
            row_count += 1
        if markers_found:
            regions.append(
                FillableRegion(
                    region_id=str(uuid4()),
                    label=ws.title,
                    sheet_or_slide=ws.title,
                    region_type="table",
                    headers=headers,
                    row_count=row_count,
                )
            )
    wb.close()
    return regions


def _detect_pptx_regions(file_bytes: bytes) -> list[FillableRegion]:
    """Scan a PPTX presentation for shapes containing {{marker}} placeholders."""
    try:
        from pptx import Presentation
    except ImportError:
        return []

    prs = Presentation(io.BytesIO(file_bytes))
    regions: list[FillableRegion] = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_label = f"Slide {slide_idx + 1}"
        for shape in slide.shapes:
            markers: set[str] = set()
            headers: list[str] = []
            region_type = "text"
            row_count = 0

            if shape.has_table:
                region_type = "table"
                table = shape.table
                # First row as headers
                if len(table.rows) > 0:
                    headers = [cell.text for cell in table.rows[0].cells]
                for row_idx, row in enumerate(table.rows):
                    for cell in row.cells:
                        markers.update(MARKER_PATTERN.findall(cell.text))
                    if row_idx > 0:
                        row_count += 1

            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        markers.update(MARKER_PATTERN.findall(run.text))

            if markers:
                regions.append(
                    FillableRegion(
                        region_id=str(uuid4()),
                        label=f"{slide_label} - {shape.name}",
                        sheet_or_slide=slide_label,
                        region_type=region_type,
                        headers=headers,
                        row_count=row_count,
                    )
                )
    return regions


# ======================================================================
# Fill helpers
# ======================================================================


def _replace_markers_in_string(text: str, fills: dict) -> str:
    """Replace all {{key}} markers in a string with fill values."""
    def _replacer(match: re.Match) -> str:
        key = match.group(1)
        return str(fills.get(key, match.group(0)))
    return MARKER_PATTERN.sub(_replacer, text)


def _fill_xlsx(template_bytes: bytes, fills: dict, regions: list[dict]) -> bytes:
    """Apply fills to an XLSX workbook and return the result bytes."""
    wb = load_workbook(io.BytesIO(template_bytes))

    # Build a lookup of region data by sheet name
    region_data: dict[str, list[dict]] = {}
    for region in regions:
        sheet_name = region.get("sheet_or_slide", "")
        if sheet_name in fills:
            region_data[sheet_name] = fills[sheet_name]

    for ws in wb.worksheets:
        # If fills contain row data for this sheet, write them.
        # Only process list values (keyed by sheet title); skip dict values
        # (keyed by region_id for frontend consumption).
        sheet_fills = fills.get(ws.title)
        if isinstance(sheet_fills, list):
            # Write rows starting at row 2
            for row_idx, row_data in enumerate(sheet_fills, start=2):
                if isinstance(row_data, dict):
                    # Map by header names from row 1
                    headers = [str(c.value) if c.value else "" for c in ws[1]]
                    for col_idx, header in enumerate(headers, start=1):
                        if header in row_data:
                            ws.cell(row=row_idx, column=col_idx, value=row_data[header])
                elif isinstance(row_data, list):
                    for col_idx, val in enumerate(row_data, start=1):
                        ws.cell(row=row_idx, column=col_idx, value=val)

        # Replace any remaining {{marker}} placeholders in all cells
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str):
                    new_val = _replace_markers_in_string(cell.value, fills)
                    if new_val != cell.value:
                        cell.value = new_val

    buf = io.BytesIO()
    wb.save(buf)
    wb.close()
    return buf.getvalue()


def _fill_pptx(template_bytes: bytes, fills: dict, regions: list[dict]) -> bytes:
    """Apply fills to a PPTX presentation and return the result bytes."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PPTX report generation")

    prs = Presentation(io.BytesIO(template_bytes))

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if MARKER_PATTERN.search(cell.text):
                            # Replace in each run to preserve formatting
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = _replace_markers_in_string(run.text, fills)

            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if MARKER_PATTERN.search(run.text):
                            run.text = _replace_markers_in_string(run.text, fills)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
